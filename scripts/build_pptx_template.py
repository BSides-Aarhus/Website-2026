"""Build a BSides Aarhus 2026 PowerPoint template matching the website design.

Generates a 16:9 .pptx with:
  - Title / cover slide
  - Agenda slide (used before talks and between talks)
  - Section divider
  - Content slide (title + body)
  - Speaker bio slide

Design tokens come from static/css/main.css.
"""
import math
import random
import re
import subprocess
import uuid
from io import BytesIO
from pathlib import Path

import yaml

from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from pptx.util import Inches, Pt, Emu
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "static" / "images" / "logo.png"
DISCORD_QR_SVG = ROOT / "static" / "images" / "discord-qr-code.svg"
SPEAKERS_DIR = ROOT / "content" / "speakers"
SESSIONS_DIR = ROOT / "content" / "sessions"
PHOTOS_DIR = ROOT / "static" / "images" / "speakers"
SCHEDULE_YAML = ROOT / "data" / "schedule.yaml"
SPONSORS_YAML = ROOT / "data" / "sponsors.yaml"
SPONSORS_DIR = ROOT / "static" / "images" / "sponsors"
OUT_DIR = ROOT / "Plans" / "template-output"
OUT = OUT_DIR / "bsides-aarhus-2026-template.pptx"
BG_IMG = OUT_DIR / "_bg.png"
BG_IMG_ELEV = OUT_DIR / "_bg_elev.png"
DISCORD_QR_PNG = OUT_DIR / "_discord_qr.png"

# --- Design tokens (from main.css) -------------------------------------------------
BG          = RGBColor(0x0A, 0x0A, 0x0F)
BG_ELEV     = RGBColor(0x12, 0x12, 0x1A)
BG_CARD     = RGBColor(0x1A, 0x1A, 0x25)
TEXT        = RGBColor(0xE8, 0xE8, 0xED)
TEXT_MUTED  = RGBColor(0x88, 0x88, 0xA0)
ACCENT      = RGBColor(0x4A, 0xB8, 0xD2)
ACCENT_LT   = RGBColor(0x6D, 0xD0, 0xE7)
BORDER      = RGBColor(0x2A, 0x2A, 0x3A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Inter"
FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

# --- Slide geometry (16:9) ---------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- Background image (white dots + top-center cyan glow) --------------------------
def build_background(path: Path, *, base=(0x0A, 0x0A, 0x0F), width=2560, height=1440,
                     dot_count=70, glow_strength=0.22):
    """Render a static version of the website's particle + radial-gradient background.

    Mirrors static/js/main.js (particles, connecting lines) and main.css
    (radial-gradient ellipse at 50% 0%). Dots and lines are rendered white per
    request; the top-center glow uses the cyan accent.
    """
    rng = random.Random(20260620)
    img = Image.new("RGB", (width, height), base)

    # Top-center cyan radial glow — drawn as a soft blurred ellipse.
    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gx, gy = width // 2, 0
    rx, ry = int(width * 0.6), int(height * 0.7)
    gd.ellipse((gx - rx, gy - ry, gx + rx, gy + ry),
               fill=(0x4A, 0xB8, 0xD2, int(255 * glow_strength)))
    glow = glow.filter(ImageFilter.GaussianBlur(radius=180))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")

    # Particles
    pts = []
    for _ in range(dot_count):
        pts.append((
            rng.uniform(0, width),
            rng.uniform(0, height),
            rng.uniform(2, 5),                 # radius in px
            rng.uniform(0.35, 0.85),           # alpha
        ))

    # Connecting lines (faint, white) — matches main.js connectDist behavior.
    connect = (width / 1920) * 150
    lines = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    ld = ImageDraw.Draw(lines)
    for i, (x1, y1, _, _) in enumerate(pts):
        for x2, y2, _, _ in pts[i + 1:]:
            d = math.hypot(x1 - x2, y1 - y2)
            if d < connect:
                a = int(255 * (1 - d / connect) * 0.12)
                ld.line((x1, y1, x2, y2), fill=(255, 255, 255, a), width=1)
    img = Image.alpha_composite(img.convert("RGBA"), lines).convert("RGB")

    # Dots (white, with a soft halo via blur+composite).
    dots = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    dd = ImageDraw.Draw(dots)
    for x, y, r, a in pts:
        dd.ellipse((x - r, y - r, x + r, y + r),
                   fill=(255, 255, 255, int(255 * a)))
    halo = dots.filter(ImageFilter.GaussianBlur(radius=6))
    img = Image.alpha_composite(img.convert("RGBA"), halo)
    img = Image.alpha_composite(img, dots).convert("RGB")

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG", optimize=True)
    return path


def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def paint_background(slide, rgb=BG, image: Path | None = None):
    if image is not None and image.exists():
        slide.shapes.add_picture(str(image), 0, 0, width=SLIDE_W, height=SLIDE_H)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        set_solid(bg, rgb)
        bg.shadow.inherit = False


def accent_bar(slide, top=Inches(0.55), left=Inches(0.6), width=Inches(0.9), height=Emu(38100)):
    """Thin cyan accent bar — used as a visual signature on most slides."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_solid(bar, ACCENT)


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=TEXT, font=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def footer(slide, page_label=""):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "BSides Aarhus 2026  ·  June 20  ·  INCUBA Next, Aarhus",
             size=10, color=TEXT_MUTED, font=FONT_MONO)
    if page_label:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
                 page_label, size=10, color=TEXT_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.RIGHT)


def add_logo(slide, left, top, height=Inches(0.6)):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), left, top, height=height)


def rasterize_discord_qr(width_px: int = 600) -> Path | None:
    """Render the Discord QR SVG to PNG via rsvg-convert. Cached in OUT_DIR."""
    if not DISCORD_QR_SVG.exists():
        return None
    import subprocess
    try:
        subprocess.run(
            ["rsvg-convert", "-w", str(width_px),
             str(DISCORD_QR_SVG), "-o", str(DISCORD_QR_PNG)],
            check=True, capture_output=True,
        )
        return DISCORD_QR_PNG
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"  warn: could not rasterize Discord QR: {e}")
        return None


def add_discord_qr(slide, *, left, top, size=Inches(1.2)):
    """Place the Discord QR (white tile + caption) at (left, top) on a slide."""
    qr_png = rasterize_discord_qr()
    if qr_png is None:
        return

    # White tile behind the QR so the dark dots scan reliably.
    pad = Inches(0.08)
    tile = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left - pad, top - pad, size + pad * 2, size + pad * 2,
    )
    tile.adjustments[0] = 0.08
    set_solid(tile, WHITE)

    slide.shapes.add_picture(str(qr_png), left, top, width=size, height=size)

    gap = Inches(0.35)
    caption_w = Inches(1.8)
    caption_left = left - caption_w - gap
    add_text(slide, caption_left, top + Inches(0.05), caption_w, Inches(0.3),
             "JOIN THE CONVERSATION", size=9, bold=True, color=ACCENT,
             font=FONT_MONO, align=PP_ALIGN.RIGHT)
    add_text(slide, caption_left, top + Inches(0.35), caption_w, Inches(0.9),
             "Scan to open the BSides Aarhus channel — invite link at bsidesaarhus.dk",
             size=8, color=TEXT_MUTED, font=FONT_BODY, align=PP_ALIGN.RIGHT)


# --- Slide builders ----------------------------------------------------------------
def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, image=BG_IMG)

    add_logo(slide, Inches(0.6), Inches(0.55), height=Inches(1.0))

    add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
             "SECURITY CONFERENCE  ·  AARHUS",
             size=14, color=ACCENT, font=FONT_MONO, bold=True)

    add_text(slide, Inches(0.6), Inches(2.7), Inches(12), Inches(1.6),
             "BSides Aarhus 2026",
             size=72, bold=True, color=TEXT, font=FONT_HEAD)

    add_text(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6),
             "Talk title goes here",
             size=32, color=ACCENT_LT, font=FONT_HEAD)

    add_text(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.45),
             "Speaker Name  ·  Affiliation",
             size=20, color=TEXT_MUTED, font=FONT_BODY)

    # Bottom meta strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, Inches(6.8), SLIDE_W, Emu(38100))
    set_solid(strip, ACCENT)
    add_text(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
             "JUNE 20, 2026  ·  INCUBA NEXT, KATRINEBJERG  ·  BSIDESAARHUS.DK",
             size=11, color=TEXT_MUTED, font=FONT_MONO)


# --- Schedule / session data (drives the agenda) -----------------------------------
ROOM_TO_TRACK = {"Store aud": 1, "Lille aud": 2}


def _read_frontmatter(md_path: Path) -> dict:
    """Return the YAML frontmatter of a markdown file as a dict (or {})."""
    text = md_path.read_text(encoding="utf-8")
    m = _FRONTMATTER_RE.match(text)
    if not m:
        return {}
    return yaml.safe_load(m.group(1)) or {}


def load_sessions() -> dict:
    """Map session slug -> {title, time, room, speakers} from content/sessions/*.md."""
    sessions = {}
    for p in SESSIONS_DIR.glob("*.md"):
        if p.stem.endswith(".da") or p.stem.startswith("_"):
            continue
        fm = _read_frontmatter(p)
        if not fm:
            continue
        sessions[p.stem] = {
            "title": str(fm.get("title", "")).strip(),
            "room": str(fm.get("room", "")).strip(),
            "speakers": fm.get("speakers", []) or [],
        }
    return sessions


def speaker_name_map() -> dict:
    """Map speaker slug -> display name (title) from content/speakers/*.md."""
    names = {}
    for p in SPEAKERS_DIR.glob("*.md"):
        if p.stem.endswith(".da") or p.stem.startswith("_"):
            continue
        fm = _read_frontmatter(p)
        if fm.get("title"):
            names[p.stem] = str(fm["title"]).strip()
    return names


def session_track_map(sessions: dict) -> dict:
    """Map speaker slug -> track number, derived from each session's room."""
    slug_track = {}
    for sess in sessions.values():
        track = ROOM_TO_TRACK.get(sess["room"])
        if not track:
            continue
        for sp in sess["speakers"]:
            slug_track[sp] = track
    return slug_track


def _fmt_time(value) -> str:
    """Normalize a schedule time to HH:MM (PyYAML may read 09:00 as sexagesimal int)."""
    if isinstance(value, int):  # YAML 1.1 base-60: 09:00 -> 540
        return f"{value // 60:02d}:{value % 60:02d}"
    return str(value).strip()


def load_agenda_rows() -> list:
    """Build agenda rows from data/schedule.yaml + content/sessions/*.md.

    Returns tuples (time, kind, t1_title, t1_speaker, t2_title, t2_speaker) where
    kind is "talk" (two-track) or "shared" (full-width break/plenary).
    """
    data = yaml.safe_load(SCHEDULE_YAML.read_text(encoding="utf-8")) or {}
    sessions = load_sessions()
    names = speaker_name_map()
    rows = []

    def resolve(slug):
        sess = sessions.get(slug, {})
        title = sess.get("title", slug or "")
        speakers = sess.get("speakers", [])
        speaker = "  ·  ".join(names.get(s, s) for s in speakers)
        return title, speaker

    for slot in data.get("timeSlots", []):
        t = _fmt_time(slot.get("time", ""))
        if slot.get("type") == "parallel":
            t1, s1 = resolve(slot.get("track1"))
            t2, s2 = resolve(slot.get("track2"))
            rows.append((t, "talk", t1, s1, t2, s2))
        else:  # break / plenary -> shared full-width row
            title = (slot.get("title") or {}).get("en", "")
            rows.append((t, "shared", title, "", "", ""))
    return rows


def build_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, image=BG_IMG)
    add_logo(slide, Inches(0.6), Inches(0.35), height=Inches(0.72))
    accent_bar(slide, top=Inches(1.1), left=Inches(0.6), width=Inches(0.7))

    add_text(slide, Inches(0.6), Inches(1.22), Inches(12), Inches(0.7),
             "Agenda", size=34, bold=True, color=TEXT, font=FONT_HEAD)
    add_text(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(0.35),
             "JUNE 20, 2026  ·  INCUBA NEXT, KATRINEBJERG",
             size=11, color=TEXT_MUTED, font=FONT_MONO)

    # Discord QR — top-right corner, caption to the left of the QR.
    add_discord_qr(slide, left=Inches(11.03), top=Inches(0.55), size=Inches(1.7))

    # Two-track schedule, loaded live from data/schedule.yaml + content/sessions/*.md.
    # (time, kind, track1_title, track1_speaker, track2_title, track2_speaker)
    # kind: "talk" | "shared" — shared rows span both columns (breaks/plenary).
    rows = load_agenda_rows()

    # Layout: time column | Track 1 | Track 2
    left_margin = Inches(0.6)
    time_w = Inches(1.05)
    gap = Inches(0.15)
    track_w = Inches((13.333 - 0.6 - 0.6 - 1.05 - 0.15 - 0.15) / 2)  # ~5.4"
    t1_left = left_margin + time_w + gap
    t2_left = t1_left + track_w + gap

    header_top = Inches(2.35)
    header_h = Inches(0.35)
    top = Inches(2.78)
    row_h = Inches(0.48)

    # Single ~15% opacity panel behind the entire agenda so dots stay visible.
    # Added BEFORE the headers/rows so it sits underneath them in z-order.
    panel_top = Inches(2.30)
    panel_h = Inches(0.50) + row_h * len(rows)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, panel_top,
        Inches(13.333) - left_margin * 2, panel_h)
    panel.adjustments[0] = 0.02
    panel.fill.solid()
    panel.fill.fore_color.rgb = BG_CARD
    panel.line.fill.background()
    from pptx.oxml.ns import qn
    sppr = panel.fill._xPr
    solidFill = sppr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    alpha = srgb.makeelement(qn('a:alpha'), {'val': '15000'})
    srgb.append(alpha)

    add_text(slide, t1_left, header_top, track_w, header_h,
             "TRACK 1  ·  STORE AUD", size=11, bold=True, color=ACCENT,
             font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(slide, t2_left, header_top, track_w, header_h,
             "TRACK 2  ·  LILLE AUD", size=11, bold=True, color=ACCENT,
             font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    for i, (t, kind, t1, s1, t2, s2) in enumerate(rows):
        y = top + row_h * i

        # Time column
        add_text(slide, left_margin, y, time_w, row_h,
                 t, size=14, color=ACCENT, font=FONT_MONO, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

        if kind == "shared":
            add_text(slide, t1_left, y,
                     Inches(13.333) - t1_left - left_margin, row_h,
                     t1, size=14, color=TEXT_MUTED, font=FONT_BODY, bold=False,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            continue

        # Track 1
        _talk_cell(slide, t1_left, y, track_w, row_h, t1, s1)
        # Track 2
        _talk_cell(slide, t2_left, y, track_w, row_h, t2, s2)

    footer(slide, "Agenda")


def _talk_cell(slide, left, top, width, height, title, speaker):
    """Render a two-line talk cell: title (top) + speaker (bottom, muted)."""
    box = slide.shapes.add_textbox(left + Inches(0.15), top, width - Inches(0.3), height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT_BODY
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = TEXT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = speaker
    r2.font.name = FONT_MONO
    r2.font.size = Pt(9)
    r2.font.color.rgb = TEXT_MUTED


def build_section_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, BG_ELEV, image=BG_IMG_ELEV)

    # Large vertical accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(2.2), Inches(0.12), Inches(3.0))
    set_solid(bar, ACCENT)

    add_text(slide, Inches(1.0), Inches(2.1), Inches(11), Inches(0.6),
             "SECTION", size=16, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(slide, Inches(1.0), Inches(2.7), Inches(11.5), Inches(2.0),
             "Section title goes here",
             size=60, bold=True, color=TEXT, font=FONT_HEAD)
    add_text(slide, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.6),
             "Optional one-line description for this section.",
             size=22, color=TEXT_MUTED, font=FONT_BODY)

    footer(slide)


def build_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, image=BG_IMG)
    add_logo(slide, Inches(0.6), Inches(0.5), height=Inches(0.72))
    accent_bar(slide, top=Inches(1.3), left=Inches(0.6), width=Inches(0.7))

    add_text(slide, Inches(0.6), Inches(1.45), Inches(12), Inches(0.8),
             "Slide title", size=36, bold=True, color=TEXT, font=FONT_HEAD)

    bullets = [
        "Primary point — short, declarative, one idea per line",
        "Supporting detail with concrete example",
        "Evidence, data, or reference",
        "Takeaway or call to action",
    ]
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5),
                                   Inches(12.1), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = "▍ "
        r1.font.name = FONT_MONO
        r1.font.size = Pt(20)
        r1.font.color.rgb = ACCENT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = FONT_BODY
        r2.font.size = Pt(22)
        r2.font.color.rgb = TEXT

    footer(slide, "Content")


def build_speaker_bio(prs, *, name="Speaker Name", role="Role  ·  Company",
                      bio=("Short bio. Two or three sentences about background, "
                           "specialization, and what the audience will take away from the "
                           "talk. Keep it conversational and specific."),
                      handle="@handle  ·  speaker.example.com",
                      photo: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, image=BG_IMG)
    add_logo(slide, Inches(0.6), Inches(0.5), height=Inches(0.72))

    # Photo card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), Inches(1.6), Inches(4.2), Inches(4.8))
    set_solid(card, BG_CARD)
    card.adjustments[0] = 0.04

    if photo and photo.exists():
        # Center-cropped square photo inside the card (EXIF-normalized).
        img = Image.open(photo)
        img = ImageOps.exif_transpose(img)
        w, h = img.size
        side = min(w, h)
        left = (w - side) // 2
        top = (h - side) // 2
        img = img.crop((left, top, left + side, top + side))
        if img.mode != "RGB":
            img = img.convert("RGB")
        buf = BytesIO()
        img.save(buf, "JPEG", quality=88)
        buf.seek(0)
        size = Inches(3.4)
        photo_left = Inches(0.6) + (Inches(4.2) - size) / 2
        photo_top = Inches(1.6) + (Inches(4.8) - size) / 2
        slide.shapes.add_picture(buf, photo_left, photo_top, width=size, height=size)
    else:
        add_text(slide, Inches(0.6), Inches(3.7), Inches(4.2), Inches(0.6),
                 "[ Photo ]", size=16, color=TEXT_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)

    accent_bar(slide, top=Inches(1.5), left=Inches(5.2), width=Inches(0.7))
    add_text(slide, Inches(5.2), Inches(1.65), Inches(7.5), Inches(0.4),
             "SPEAKER", size=14, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(slide, Inches(5.2), Inches(2.1), Inches(7.5), Inches(0.9),
             name, size=40, bold=True, color=TEXT, font=FONT_HEAD)
    add_text(slide, Inches(5.2), Inches(3.05), Inches(7.5), Inches(0.5),
             role, size=20, color=ACCENT_LT, font=FONT_BODY)
    add_text(slide, Inches(5.2), Inches(3.8), Inches(7.5), Inches(2.3),
             bio, size=16, color=TEXT_MUTED, font=FONT_BODY)
    add_text(slide, Inches(5.2), Inches(6.2), Inches(7.5), Inches(0.4),
             handle, size=14, color=ACCENT, font=FONT_MONO)

    footer(slide, "Speaker")


# --- Speaker data ------------------------------------------------------------------
_FRONTMATTER_RE = re.compile(r"^---\n(.*?)\n---\n(.*)$", re.DOTALL)


def parse_speaker(md_path: Path):
    text = md_path.read_text(encoding="utf-8")
    m = _FRONTMATTER_RE.match(text)
    if not m:
        return None
    fm_raw, body = m.group(1), m.group(2).strip()
    fm, links, in_links = {}, {}, False
    for line in fm_raw.splitlines():
        if not line.strip():
            continue
        if line.startswith("links:"):
            in_links = True
            continue
        if in_links and line.startswith("  "):
            k, _, v = line.strip().partition(":")
            links[k.strip()] = v.strip().strip('"').strip("'")
            continue
        in_links = False
        k, _, v = line.partition(":")
        fm[k.strip()] = v.strip().strip('"').strip("'")
    paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()]
    fm["bio"] = paragraphs[0] if paragraphs else ""
    fm["links"] = links
    fm["slug"] = md_path.stem
    return fm


def find_photo(slug: str):
    for ext in ("jpg", "jpeg", "png"):
        p = PHOTOS_DIR / f"{slug}.{ext}"
        if p.exists():
            return p
    return None


def build_handle_line(links: dict) -> str:
    parts = []
    if "linkedin" in links:
        m = re.search(r"/in/([^/]+)/?", links["linkedin"])
        parts.append(f"linkedin.com/in/{m.group(1)}" if m else links["linkedin"])
    if "website" in links:
        parts.append(re.sub(r"^https?://", "", links["website"]).rstrip("/"))
    if "twitter" in links:
        m = re.search(r"twitter\.com/([^/]+)", links["twitter"])
        parts.append(f"@{m.group(1)}" if m else links["twitter"])
    if "mastodon" in links:
        parts.append(links["mastodon"])
    return "  ·  ".join(parts)


def load_speakers():
    files = sorted(
        p for p in SPEAKERS_DIR.glob("*.md")
        if not p.stem.endswith(".da") and not p.stem.startswith("_")
    )
    # Primary source of truth: each session's room ("Store aud"=1, "Lille aud"=2).
    slug_track = session_track_map(load_sessions())
    speakers = []
    for sf in files:
        s = parse_speaker(sf)
        if not s:
            continue
        s["photo"] = find_photo(s["slug"])
        if s["slug"] in slug_track:
            s["track"] = slug_track[s["slug"]]
        else:
            print(f"  warn: {s['slug']} has no session room — placing in Track 1")
            s["track"] = 1
        speakers.append(s)
    # Track 1 alpha, then Track 2 alpha
    return (sorted([s for s in speakers if s["track"] == 1], key=lambda x: x["title"])
          + sorted([s for s in speakers if s["track"] == 2], key=lambda x: x["title"]))


# --- PowerPoint sections (Template / Track 1 / Track 2) ----------------------------
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def apply_fade_transition(prs, duration_ms: int = 700):
    """Add a fade transition to every slide so they don't switch instantly."""
    MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    for slide in prs.slides:
        sld = slide._element
        for existing in sld.findall(qn("p:transition")):
            sld.remove(existing)
        for existing in sld.findall(f"{{{MC_NS}}}AlternateContent"):
            sld.remove(existing)
        alt = etree.SubElement(sld, f"{{{MC_NS}}}AlternateContent",
                               nsmap={"mc": MC_NS})
        choice = etree.SubElement(alt, f"{{{MC_NS}}}Choice",
                                  nsmap={"p14": P14_NS}, Requires="p14")
        t1 = etree.SubElement(choice, qn("p:transition"),
                              attrib={"spd": "med", f"{{{P14_NS}}}dur": str(duration_ms)})
        etree.SubElement(t1, qn("p:fade"))
        fb = etree.SubElement(alt, f"{{{MC_NS}}}Fallback")
        t2 = etree.SubElement(fb, qn("p:transition"), spd="med")
        etree.SubElement(t2, qn("p:fade"))


def add_sections(prs, template_count: int, speakers: list):
    """Inject Template / Track 1 / Track 2 sections into the presentation XML."""
    pres = prs.part._element  # <p:presentation>

    # Collect <p:sldId> ids in document order.
    sld_ids = [int(e.get("id")) for e in pres.findall(qn("p:sldIdLst") + "/" + qn("p:sldId"))]
    template_ids = sld_ids[:template_count]
    speaker_ids = sld_ids[template_count:]
    track1 = [sid for sid, sp in zip(speaker_ids, speakers) if sp["track"] == 1]
    track2 = [sid for sid, sp in zip(speaker_ids, speakers) if sp["track"] == 2]

    # Find or create <p:extLst> as the last child of <p:presentation>.
    ext_lst = pres.find(qn("p:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(pres, qn("p:extLst"))

    # Remove any existing section ext to make this idempotent.
    SECTION_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
    for ext in ext_lst.findall(qn("p:ext")):
        if ext.get("uri") == SECTION_URI:
            ext_lst.remove(ext)

    ext = etree.SubElement(ext_lst, qn("p:ext"), uri=SECTION_URI)
    section_lst = etree.SubElement(ext, f"{{{P14_NS}}}sectionLst", nsmap={"p14": P14_NS})

    def add_section(name, ids):
        sec = etree.SubElement(section_lst, f"{{{P14_NS}}}section",
                               name=name, id="{" + str(uuid.uuid4()).upper() + "}")
        sld_id_lst = etree.SubElement(sec, f"{{{P14_NS}}}sldIdLst")
        for i in ids:
            etree.SubElement(sld_id_lst, f"{{{P14_NS}}}sldId", id=str(i))

    add_section("Template", template_ids)
    add_section("Track 1", track1)
    add_section("Track 2", track2)


# --- Sponsors ----------------------------------------------------------------------
def load_sponsor_tiers() -> list:
    """Return [{label, sponsors:[{name,url,logo,darkBg}]}] for non-empty tiers."""
    data = yaml.safe_load(SPONSORS_YAML.read_text(encoding="utf-8")) or {}
    tiers = []
    for tier in data.get("tiers", []):
        sponsors = tier.get("sponsors") or []
        if not sponsors:
            continue  # skip empty tiers (e.g. Community Food Friends)
        tiers.append({
            "label": (tier.get("name") or {}).get("en", ""),
            "sponsors": [{
                "name": sp.get("name", ""),
                "logo": sp.get("logo", ""),
                "darkBg": bool(sp.get("darkBg", False)),
            } for sp in sponsors],
        })
    return tiers


def resolve_logo(logo_rel: str) -> Path | None:
    """Resolve a sponsors.yaml logo path to a raster file (rasterizing SVG via rsvg)."""
    if not logo_rel:
        return None
    src = SPONSORS_DIR / Path(logo_rel).name
    if not src.exists():
        print(f"  warn: sponsor logo not found: {src.name}")
        return None
    if src.suffix.lower() != ".svg":
        return src
    out = OUT_DIR / f"_sponsor_{src.stem}.png"
    try:
        subprocess.run(["rsvg-convert", "-h", "600", str(src), "-o", str(out)],
                       check=True, capture_output=True)
        return out
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"  warn: could not rasterize {src.name}: {e}")
        return None


def build_sponsors(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, image=BG_IMG)
    add_logo(slide, Inches(0.6), Inches(0.35), height=Inches(0.72))
    accent_bar(slide, top=Inches(1.1), left=Inches(0.6), width=Inches(0.7))

    add_text(slide, Inches(0.6), Inches(1.22), Inches(12), Inches(0.7),
             "Thank You to Our Sponsors", size=34, bold=True, color=TEXT, font=FONT_HEAD)

    # Flatten every sponsor across tiers — the wall shows all logos equally.
    sponsors = []
    for tier in load_sponsor_tiers():
        for sp in tier["sponsors"]:
            path = resolve_logo(sp["logo"])
            if path is not None:
                with Image.open(path) as im:
                    aspect = im.width / im.height if im.height else 3.0
                sponsors.append({"sp": sp, "path": path, "aspect": aspect})
    if not sponsors:
        footer(slide, "Sponsors")
        return

    # --- Equal-height logo wall ----------------------------------------------------
    # Upright (no rotation), professional. Every logo is drawn at the SAME height for
    # uniform prominence and sized as large as fits; rows are justified to the content
    # width. Each logo sits on a subtle uniform card; dark logos get a white plate.
    L, R = 0.55, 12.78
    TOP, BOT = 2.02, 7.02           # below the title, above the footer
    content_w = R - L
    region_h = BOT - TOP
    HGAP, VGAP = 0.36, 0.3          # min gap between logos in a row / between rows
    PAD = 0.14                      # card padding around each logo
    MAX_GAP = 1.2                   # wider justify gap than this → center the row

    def pack(h):
        """First-fit-decreasing: widest logos first, each into the first row with room.
        Pairing wide wordmarks with small marks keeps rows full and avoids a lone logo."""
        rows = []
        for s in sorted(sponsors, key=lambda s: s["aspect"], reverse=True):
            w = h * s["aspect"]
            for row in rows:
                if row["w"] + HGAP + w <= content_w:
                    row["items"].append(s)
                    row["w"] += HGAP + w
                    break
            else:
                rows.append({"items": [s], "w": w})
        return [r["items"] for r in rows]

    # Largest uniform height: each row fits content_w (by construction) and, with the
    # row gaps, the whole stack fits region_h. A single logo must also fit one row.
    widest = max(s["aspect"] for s in sponsors)
    H, rows = 0.6, pack(0.6)
    for i in range(1, 120):
        h = 0.6 + 0.01 * i
        if h * widest > content_w:
            break
        cand = pack(h)
        if len(cand) * h + (len(cand) - 1) * VGAP <= region_h:
            H, rows = h, cand
        else:
            break

    total_h = len(rows) * H + (len(rows) - 1) * VGAP
    y = TOP + (region_h - total_h) / 2
    for row in rows:
        widths = [H * s["aspect"] for s in row]
        nat, n = sum(widths), len(row)
        if n > 1 and (content_w - nat) / (n - 1) <= MAX_GAP:
            gap, x = (content_w - nat) / (n - 1), L       # justify across full width
        else:
            gap, x = HGAP, L + (content_w - (nat + HGAP * (n - 1))) / 2   # center row

        for s, w in zip(row, widths):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x - PAD), Inches(y - PAD), Inches(w + 2 * PAD), Inches(H + 2 * PAD))
            card.adjustments[0] = 0.10
            card.shadow.inherit = False
            card.fill.solid()
            card.fill.fore_color.rgb = BG_CARD
            card.line.color.rgb = BORDER
            card.line.width = Pt(0.75)
            csf = card.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
            csf.append(csf.makeelement(qn('a:alpha'), {'val': '15000'}))

            if s["sp"]["darkBg"]:                         # white plate for dark artwork
                pp = 0.06
                plate = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x - pp), Inches(y - pp), Inches(w + 2 * pp), Inches(H + 2 * pp))
                plate.adjustments[0] = 0.10
                set_solid(plate, WHITE)
                plate.shadow.inherit = False

            slide.shapes.add_picture(str(s["path"]), Inches(x), Inches(y),
                                     width=Inches(w), height=Inches(H))
            x += w + gap
        y += H + VGAP

    footer(slide, "Sponsors")


def main():
    # Pre-render the dotted backgrounds (base + elevated tone for dividers).
    build_background(BG_IMG, base=(0x0A, 0x0A, 0x0F))
    build_background(BG_IMG_ELEV, base=(0x12, 0x12, 0x1A), glow_strength=0.18)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Template section (5 slides) ---
    build_cover(prs)
    build_agenda(prs)
    build_section_divider(prs)
    build_content(prs)
    build_speaker_bio(prs)  # generic placeholder template
    build_sponsors(prs)     # sponsors, grouped by tier from data/sponsors.yaml
    template_count = len(prs.slides)

    # --- Speaker slides, grouped Track 1 then Track 2 ---
    speakers = load_speakers()
    print(f"Adding {len(speakers)} speaker slides "
          f"(Track 1: {sum(1 for s in speakers if s['track']==1)}, "
          f"Track 2: {sum(1 for s in speakers if s['track']==2)})")
    for s in speakers:
        build_speaker_bio(
            prs,
            name=s["title"],
            role=s.get("tagline", ""),
            bio=s["bio"],
            handle=build_handle_line(s["links"]),
            photo=s["photo"],
        )

    # --- Sections ---
    add_sections(prs, template_count, speakers)

    # --- Transitions (fade between slides) ---
    apply_fade_transition(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
